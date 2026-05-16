from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Цветовая схема в стиле приложения
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x33, 0x33, 0x66)
TEXT_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_GRAY = RGBColor(0xAA, 0xAA, 0xCC)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_ORANGE = RGBColor(0xFF, 0xA5, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_LIGHT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.level = 0
        p.space_after = Pt(6)
    return tf

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

# ============================================================
# СЛАЙД 1: Титульный
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
set_slide_bg(slide1, BG_DARK)

add_text_box(slide1, 1.5, 1.5, 10, 1.2, "🌤 Погода Омск", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 2.8, 10, 0.8, "Умный помощник по выбору одежды", font_size=28, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_accent_line(slide1, 4.5, 3.8, 4)
add_text_box(slide1, 1.5, 4.2, 10, 0.8, "Python • Tkinter • wttr.in API", font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 5.8, 10, 0.6, "Нажмите Пробел или → для продолжения", font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# СЛАЙД 2: Обзор проекта
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_text_box(slide2, 0.8, 0.5, 6, 0.8, "📋 О проекте", font_size=36, color=WHITE, bold=True)
add_accent_line(slide2, 0.8, 1.3, 3)

add_text_box(slide2, 0.8, 1.8, 5.5, 1.5, 
    "Десктопное приложение, которое\n"
    "анализирует погоду в Омске\n"
    "и рекомендует подходящую одежду.",
    font_size=20, color=TEXT_LIGHT)

add_bullet_list(slide2, 0.8, 3.5, 5.5, 3.5, [
    "🎯 Цель: автоматический подбор гардероба",
    "🐍 Язык: Python 3",
    "🖥 Интерфейс: Tkinter (тёмная тема)",
    "🌐 Данные: wttr.in API (JSON)",
    "📦 Размер окна: 280×350 px",
], font_size=18)

# Правая колонка — архитектура
add_text_box(slide2, 7.5, 0.5, 5, 0.8, "⚙ Как это работает", font_size=28, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide2, 7.5, 1.5, 5, 5, [
    "1.  Запрос к wttr.in",
    "    └─ формат j1 (JSON)",
    "    └─ язык: русский",
    "",
    "2.  Извлечение данных:",
    "    └─ Температура (°C)",
    "    └─ Скорость ветра",
    "    └─ Влажность",
    "    └─ Описание погоды",
    "",
    "3.  Логика подбора одежды",
    "",
    "4.  Вывод в GUI",
], font_size=16)

# ============================================================
# СЛАЙД 3: Ключевые функции
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_text_box(slide3, 0.8, 0.3, 12, 0.8, "⚡ Ключевые функции", font_size=36, color=WHITE, bold=True)
add_accent_line(slide3, 0.8, 1.1, 3)

# Блок 1
add_text_box(slide3, 0.8, 1.8, 3.5, 0.5, "🌡 Температура", font_size=22, color=ACCENT_RED, bold=True)
add_text_box(slide3, 0.8, 2.3, 3.5, 2, 
    "Отображение текущей\n"
    "температуры крупным\n"
    "шрифтом (32pt)\n"
    "с цветовым акцентом.",
    font_size=16, color=TEXT_LIGHT)

# Блок 2
add_text_box(slide3, 4.9, 1.8, 3.5, 0.5, "👔 Подбор одежды", font_size=22, color=ACCENT_ORANGE, bold=True)
add_text_box(slide3, 4.9, 2.3, 3.5, 3.5,
    "10 пороговых значений\n"
    "температуры:\n"
    "• >30°C — шорты\n"
    "• 25–30°C — футболка\n"
    "• 20–25°C — ветровка\n"
    "• ...\n"
    "• <-20°C — всё тёплое!",
    font_size=15, color=TEXT_LIGHT)

# Блок 3
add_text_box(slide3, 9, 1.8, 3.8, 0.5, "💡 Умные советы", font_size=22, color=RGBColor(0x00, 0xCC, 0x88), bold=True)
add_text_box(slide3, 9, 2.3, 3.8, 2.5,
    "Дополнительные\n"
    "рекомендации:\n"
    "• Ветер >30 км/ч\n"
    "• Влажность >80%\n"
    "• Дождь — взять зонт\n"
    "• Снег — одеться теплее",
    font_size=16, color=TEXT_LIGHT)

# Разделители между колонками
for x in [4.5, 8.6]:
    shape = slide3.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(0.02), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

# ============================================================
# СЛАЙД 4: Интерфейс
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)

add_text_box(slide4, 0.8, 0.5, 12, 0.8, "🖥 Пользовательский интерфейс", font_size=36, color=WHITE, bold=True)
add_accent_line(slide4, 0.8, 1.3, 3)

# Макет окна
add_text_box(slide4, 0.8, 2.0, 5, 0.5, "Структура окна:", font_size=24, color=TEXT_GRAY, bold=True)

elements = [
    ("🌤 Заголовок «Омск»", "20pt, белый"),
    ("🌡 Температура", "32pt, красный (#FF6B6B)"),
    ("📝 Описание погоды", "12pt, серый"),
    ("━━━ Разделитель ━━━", "линия #333366"),
    ("👔 Рекомендация одежды", "11pt, белый"),
    ("💡 Советы", "10pt, оранжевый"),
    ("🔄 Кнопка «Обновить»", "11pt, с ховер-эффектом"),
]
for i, (el, desc) in enumerate(elements):
    y = 2.8 + i * 0.6
    add_text_box(slide4, 1.2, y, 3, 0.4, el, font_size=16, color=TEXT_LIGHT)
    add_text_box(slide4, 4.5, y, 3, 0.4, desc, font_size=14, color=TEXT_GRAY)

# Параметры окна справа
add_text_box(slide4, 8.5, 2.0, 4.5, 0.5, "Параметры:", font_size=24, color=TEXT_GRAY, bold=True)
params = [
    "📐 Размер: 280×350",
    "🎨 Фон: #1A1A2E",
    "🔒 Не растягивается",
    "🔄 Автообновление при запуске",
    "⚠ Обработка ошибок сети",
]
add_bullet_list(slide4, 8.5, 2.8, 4.5, 3, params, font_size=18, color=TEXT_LIGHT)

# ============================================================
# СЛАЙД 5: Технологии и перспективы
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)

add_text_box(slide5, 0.8, 0.5, 12, 0.8, "🚀 Технологии и развитие", font_size=36, color=WHITE, bold=True)
add_accent_line(slide5, 0.8, 1.3, 3)

# Стек технологий
add_text_box(slide5, 0.8, 2.0, 5.5, 0.5, "🛠 Стек технологий", font_size=24, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide5, 0.8, 2.7, 5.5, 3.5, [
    "• Python 3 — основной язык",
    "• Tkinter — GUI (встроен в Python)",
    "• requests — HTTP-запросы к API",
    "• wttr.in — бесплатный погодный сервис",
    "• JSON — формат данных",
    "• Обработка исключений (try/except)",
], font_size=17)

# Перспективы
add_text_box(slide5, 7.5, 2.0, 5.5, 0.5, "🔮 Идеи для развития", font_size=24, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide5, 7.5, 2.7, 5.5, 4, [
    "📱 Мобильная версия (Kivy / Flet)",
    "🗺 Выбор города пользователем",
    "📊 Прогноз на несколько дней",
    "🔔 Уведомления утром",
    "🎨 Кастомизация темы",
    "📎 Интеграция с Telegram-ботом",
    "🌍 Локализация на другие языки",
], font_size=17)

# Нижняя плашка
shape = slide5.shapes.add_shape(1, Inches(0), Inches(6.6), Inches(13.333), Inches(0.9))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide5, 0.5, 6.75, 12.5, 0.6,
    "💻 Исходный код доступен для доработки | API wttr.in не требует ключа | Кроссплатформенность (Windows/Linux/macOS)",
    font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

# ============================================================
# Сохранение
# ============================================================
prs.save('WeatherApp.pptx')
print("✅ Презентация сохранена: WeatherApp.pptx")